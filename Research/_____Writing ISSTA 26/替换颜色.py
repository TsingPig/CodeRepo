from pptx import Presentation
from pptx.dml.color import RGBColor

COLOR_A = RGBColor(255, 0, 0)   # 红
COLOR_B = RGBColor(0, 170, 255) # 蓝

prs = Presentation("Fig_VRAgent.pptx")

def replace_color(shape):
    # 填充色
    if shape.fill and hasattr(shape.fill, "foreground_color"):
        fc = shape.fill.foreground_color.rgb
        if fc == COLOR_A:
            shape.fill.foreground_color.rgb = COLOR_B

    # 线条色
    if shape.line and hasattr(shape.line.color, "rgb"):
        lc = shape.line.color.rgb
        if lc == COLOR_A:
            shape.line.color.rgb = COLOR_B

    # 文本颜色
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.font.color.rgb == COLOR_A:
                    r.font.color.rgb = COLOR_B

for slide in prs.slides:
    for shape in slide.shapes:
        replace_color(shape)

prs.save("output.pptx")
